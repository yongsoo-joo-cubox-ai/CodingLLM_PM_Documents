"""
Secern AI PPT Template Helpers

오렌지 브랜딩 색상 상수 + 7가지 슬라이드 타입 헬퍼 함수.
build_presentation.py와 다른 PPT 스크립트가 import하여 재사용한다.

Usage:
    # 모듈로 사용
    from secern_template_helpers import *

    # 직접 실행 시 7장짜리 샘플 PPT 생성
    python secern_template_helpers.py
"""

from __future__ import annotations

from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Cm, Pt, Emu
from pptx.oxml.ns import qn

# ---------------------------------------------------------------------------
# Secern AI Brand Colors (Orange Branding)
# ---------------------------------------------------------------------------

PRIMARY = RGBColor(0xFF, 0x48, 0x00)       # #ff4800 - 오렌지 (헤더, CTA, 강조)
SECONDARY = RGBColor(0xB7, 0x3D, 0x0C)     # #b73d0c - 진한 오렌지 (부제목, 호버)
ACCENT = RGBColor(0xFF, 0x6B, 0x2B)        # #ff6b2b - 밝은 오렌지 (테이블 헤더, 배지)
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)       # #1a1a2e - 진한 네이비 (표지 배경, 코드 블록)
LIGHT_BG = RGBColor(0xFF, 0xF5, 0xF0)      # #fff5f0 - 연한 오렌지 (카드 배경, 테이블 짝수행)
GRAY_BG = RGBColor(0xF8, 0xF9, 0xFA)       # #f8f9fa - 연한 회색 (테이블 홀수행)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)      # #333333 - 본문 텍스트
WHITE = RGBColor(0xFF, 0xFF, 0xFF)          # #ffffff
BLACK = RGBColor(0x00, 0x00, 0x00)          # #000000

CODE_BG = DARK_BG                           # 코드 블록 배경
CODE_TEXT = RGBColor(0xE0, 0xE0, 0xE0)      # 코드 텍스트 색상
BORDER_COLOR = "CCCCCC"                     # 테이블 테두리

# ---------------------------------------------------------------------------
# Font
# ---------------------------------------------------------------------------

FONT_NAME = "Pretendard"
FONT_FALLBACK = "Segoe UI"
FONT_FALLBACK_KO = "Malgun Gothic"

# ---------------------------------------------------------------------------
# Slide Dimensions (16:9)
# ---------------------------------------------------------------------------

SLIDE_WIDTH = Cm(33.867)
SLIDE_HEIGHT = Cm(19.05)

# Layout
HEADER_HEIGHT = Cm(1.5)
FOOTER_HEIGHT = Cm(0.8)
MARGIN_LEFT = Cm(1.2)
MARGIN_RIGHT = Cm(1.2)
CONTENT_TOP = Cm(2.2)
CONTENT_WIDTH = SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT

# Font sizes
TITLE_SIZE = Pt(24)
SUBTITLE_SIZE = Pt(18)
BODY_SIZE = Pt(14)
TABLE_SIZE = Pt(12)
SMALL_SIZE = Pt(10)
FOOTER_SIZE = Pt(8)

# Global slide counter
_slide_number = 0


def reset_slide_counter():
    """슬라이드 카운터를 리셋한다."""
    global _slide_number
    _slide_number = 0


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------

def _set_shape_fill(shape, color: RGBColor) -> None:
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, left, top, width, height, color: RGBColor):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    _set_shape_fill(shape, color)
    return shape


def _set_font(run, size=BODY_SIZE, bold=False, italic=False, color=DARK_TEXT, name=None):
    font = run.font
    font.size = size
    font.bold = bold
    font.italic = italic
    font.color.rgb = color
    font.name = name or FONT_NAME


def _add_textbox(slide, left, top, width, height, text="",
                 font_size=BODY_SIZE, bold=False, color=DARK_TEXT,
                 alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                 word_wrap=True):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = word_wrap
    tf.auto_size = None

    try:
        tf.paragraphs[0].alignment = alignment
    except Exception:
        pass

    try:
        txbox.text_frame._txBody.bodyPr.set("anchor", {
            MSO_ANCHOR.TOP: "t",
            MSO_ANCHOR.MIDDLE: "ctr",
            MSO_ANCHOR.BOTTOM: "b",
        }.get(anchor, "t"))
    except Exception:
        pass

    if text:
        p = tf.paragraphs[0]
        p.alignment = alignment
        run = p.add_run()
        run.text = text
        _set_font(run, size=font_size, bold=bold, color=color)

    return txbox


def _add_bullets_to_frame(tf, bullets: list[str], font_size=BODY_SIZE,
                          color=DARK_TEXT, bold=False, level=0,
                          line_spacing=Pt(24)):
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.level = level
        p.space_after = line_spacing
        p.alignment = PP_ALIGN.LEFT

        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        buChar = pPr.makeelement(qn("a:buChar"), {"char": "\u2022"})
        for existing in pPr.findall(qn("a:buChar")):
            pPr.remove(existing)
        for existing in pPr.findall(qn("a:buNone")):
            pPr.remove(existing)
        pPr.append(buChar)

        run = p.add_run()
        run.text = bullet
        _set_font(run, size=font_size, bold=bold, color=color)


def _set_cell_fill(cell, color: RGBColor) -> None:
    tcPr = cell._tc.get_or_add_tcPr()
    solidFill = tcPr.makeelement(qn("a:solidFill"), {})
    srgbClr = solidFill.makeelement(qn("a:srgbClr"), {"val": str(color)})
    solidFill.append(srgbClr)
    for existing in tcPr.findall(qn("a:solidFill")):
        tcPr.remove(existing)
    tcPr.append(solidFill)


def _set_cell_margins(cell, margin=Cm(0.15)) -> None:
    tcPr = cell._tc.get_or_add_tcPr()
    tcPr.set("marL", str(int(margin)))
    tcPr.set("marR", str(int(margin)))
    tcPr.set("marT", str(int(Cm(0.08))))
    tcPr.set("marB", str(int(Cm(0.08))))


def _set_table_borders(table, n_rows: int, n_cols: int) -> None:
    border_width = "6350"  # 0.5pt in EMU
    for row_idx in range(n_rows):
        for col_idx in range(n_cols):
            cell = table.cell(row_idx, col_idx)
            tcPr = cell._tc.get_or_add_tcPr()
            for border_name in ["a:lnL", "a:lnR", "a:lnT", "a:lnB"]:
                ln = tcPr.makeelement(qn(border_name), {"w": border_width})
                solidFill = ln.makeelement(qn("a:solidFill"), {})
                srgbClr = solidFill.makeelement(
                    qn("a:srgbClr"), {"val": BORDER_COLOR}
                )
                solidFill.append(srgbClr)
                ln.append(solidFill)
                for existing in tcPr.findall(qn(border_name)):
                    tcPr.remove(existing)
                tcPr.append(ln)


# ---------------------------------------------------------------------------
# Header / Footer / Slide Frame
# ---------------------------------------------------------------------------

def _add_header_bar(slide, title_text: str) -> None:
    _add_rect(slide, Cm(0), Cm(0), SLIDE_WIDTH, HEADER_HEIGHT, PRIMARY)
    _add_textbox(
        slide, MARGIN_LEFT, Cm(0), CONTENT_WIDTH, HEADER_HEIGHT,
        text=title_text, font_size=TITLE_SIZE, bold=True, color=WHITE,
        alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
    )


def _add_footer(slide) -> None:
    global _slide_number
    _add_textbox(
        slide, MARGIN_LEFT, SLIDE_HEIGHT - FOOTER_HEIGHT,
        Cm(6), FOOTER_HEIGHT,
        text="Secern AI", font_size=FOOTER_SIZE, bold=False,
        color=SECONDARY, alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.BOTTOM,
    )
    _add_textbox(
        slide, SLIDE_WIDTH - MARGIN_RIGHT - Cm(3),
        SLIDE_HEIGHT - FOOTER_HEIGHT, Cm(3), FOOTER_HEIGHT,
        text=str(_slide_number), font_size=FOOTER_SIZE, bold=False,
        color=SECONDARY, alignment=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.BOTTOM,
    )


def _new_slide(prs) -> "Slide":
    global _slide_number
    _slide_number += 1
    layout = prs.slide_layouts[6]  # Blank layout
    return prs.slides.add_slide(layout)


# ---------------------------------------------------------------------------
# 7 Slide Type Helpers
# ---------------------------------------------------------------------------

def add_cover_slide(prs, title: str, subtitle: str,
                    extra_lines: list[str] | None = None) -> None:
    """표지 슬라이드. DARK_BG 배경 + 오렌지 액센트."""
    slide = _new_slide(prs)
    _add_rect(slide, Cm(0), Cm(0), SLIDE_WIDTH, SLIDE_HEIGHT, DARK_BG)

    # 오렌지 액센트 라인
    _add_rect(slide, Cm(3), Cm(8.8), Cm(6), Cm(0.08), PRIMARY)

    _add_textbox(
        slide, Cm(3), Cm(5), SLIDE_WIDTH - Cm(6), Cm(4),
        text=title, font_size=Pt(40), bold=True, color=WHITE,
        alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM,
    )
    _add_textbox(
        slide, Cm(3), Cm(9.5), SLIDE_WIDTH - Cm(6), Cm(2),
        text=subtitle, font_size=SUBTITLE_SIZE, bold=False, color=LIGHT_BG,
        alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP,
    )

    if extra_lines:
        y = Cm(13)
        for line in extra_lines:
            _add_textbox(
                slide, Cm(3), y, SLIDE_WIDTH - Cm(6), Cm(1),
                text=line, font_size=SMALL_SIZE, bold=False, color=LIGHT_BG,
                alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP,
            )
            y += Cm(0.8)

    _add_footer(slide)


def add_content_slide(prs, title: str, bullets: list[str],
                      screenshot_path: str | None = None,
                      screenshots_dir: Path | None = None) -> None:
    """일반 콘텐츠 슬라이드. 오렌지 상단 바 + 제목 + 본문."""
    slide = _new_slide(prs)
    _add_header_bar(slide, title)

    img_path = None
    if screenshot_path and screenshots_dir:
        p = screenshots_dir / screenshot_path
        if p.exists():
            img_path = str(p)

    if img_path:
        text_width = Cm(12)
        img_left = MARGIN_LEFT + text_width + Cm(0.5)
        img_width = SLIDE_WIDTH - img_left - MARGIN_RIGHT
        img_height = SLIDE_HEIGHT - CONTENT_TOP - FOOTER_HEIGHT - Cm(0.5)

        txbox = _add_textbox(
            slide, MARGIN_LEFT, CONTENT_TOP, text_width,
            SLIDE_HEIGHT - CONTENT_TOP - FOOTER_HEIGHT,
        )
        _add_bullets_to_frame(txbox.text_frame, bullets, font_size=BODY_SIZE)

        try:
            slide.shapes.add_picture(
                img_path, img_left, CONTENT_TOP, img_width, img_height,
            )
        except Exception as exc:
            print(f"    [WARN] 이미지 삽입 실패 ({screenshot_path}): {exc}")
            _add_textbox(
                slide, img_left, CONTENT_TOP + Cm(3), img_width, Cm(2),
                text=f"[Screenshot: {screenshot_path}]",
                font_size=SMALL_SIZE, color=SECONDARY,
                alignment=PP_ALIGN.CENTER,
            )
    else:
        txbox = _add_textbox(
            slide, MARGIN_LEFT, CONTENT_TOP, CONTENT_WIDTH,
            SLIDE_HEIGHT - CONTENT_TOP - FOOTER_HEIGHT,
        )
        _add_bullets_to_frame(txbox.text_frame, bullets, font_size=BODY_SIZE)

    _add_footer(slide)


def add_two_column_slide(prs, title: str, left_bullets: list[str],
                         right_bullets: list[str],
                         left_title: str = "", right_title: str = "") -> None:
    """좌우 50:50 비교 슬라이드."""
    slide = _new_slide(prs)
    _add_header_bar(slide, title)

    col_width = (CONTENT_WIDTH - Cm(1)) / 2
    content_height = SLIDE_HEIGHT - CONTENT_TOP - FOOTER_HEIGHT - Cm(0.3)

    left_top = CONTENT_TOP
    if left_title:
        _add_textbox(
            slide, MARGIN_LEFT, left_top, col_width, Cm(1),
            text=left_title, font_size=Pt(16), bold=True, color=PRIMARY,
            alignment=PP_ALIGN.LEFT,
        )
        left_top += Cm(1.2)

    txbox_l = _add_textbox(
        slide, MARGIN_LEFT, left_top, col_width,
        content_height - (left_top - CONTENT_TOP),
    )
    _add_bullets_to_frame(txbox_l.text_frame, left_bullets, font_size=Pt(13))

    right_left = MARGIN_LEFT + col_width + Cm(1)
    right_top = CONTENT_TOP
    if right_title:
        _add_textbox(
            slide, right_left, right_top, col_width, Cm(1),
            text=right_title, font_size=Pt(16), bold=True, color=PRIMARY,
            alignment=PP_ALIGN.LEFT,
        )
        right_top += Cm(1.2)

    txbox_r = _add_textbox(
        slide, right_left, right_top, col_width,
        content_height - (right_top - CONTENT_TOP),
    )
    _add_bullets_to_frame(txbox_r.text_frame, right_bullets, font_size=Pt(13))

    _add_footer(slide)


def add_image_text_slide(prs, title: str, image_path: str,
                         bullets: list[str]) -> None:
    """좌 이미지 + 우 텍스트 슬라이드."""
    slide = _new_slide(prs)
    _add_header_bar(slide, title)

    img_width = (CONTENT_WIDTH - Cm(1)) * 0.55
    text_width = CONTENT_WIDTH - img_width - Cm(1)
    content_height = SLIDE_HEIGHT - CONTENT_TOP - FOOTER_HEIGHT - Cm(0.5)

    try:
        slide.shapes.add_picture(
            image_path, MARGIN_LEFT, CONTENT_TOP, img_width, content_height,
        )
    except Exception as exc:
        print(f"    [WARN] 이미지 삽입 실패: {exc}")
        _add_rect(slide, MARGIN_LEFT, CONTENT_TOP, img_width, content_height, GRAY_BG)
        _add_textbox(
            slide, MARGIN_LEFT, CONTENT_TOP + Cm(3), img_width, Cm(2),
            text="[Image placeholder]",
            font_size=SMALL_SIZE, color=SECONDARY,
            alignment=PP_ALIGN.CENTER,
        )

    text_left = MARGIN_LEFT + img_width + Cm(1)
    txbox = _add_textbox(
        slide, text_left, CONTENT_TOP, text_width, content_height,
    )
    _add_bullets_to_frame(txbox.text_frame, bullets, font_size=BODY_SIZE)

    _add_footer(slide)


def add_table_slide(prs, title: str, headers: list[str],
                    rows: list[list[str]]) -> None:
    """오렌지 헤더 테이블 슬라이드."""
    slide = _new_slide(prs)
    _add_header_bar(slide, title)

    n_rows = len(rows) + 1
    n_cols = len(headers)

    table_top = CONTENT_TOP + Cm(0.3)
    table_height = SLIDE_HEIGHT - table_top - FOOTER_HEIGHT - Cm(0.5)

    tbl_shape = slide.shapes.add_table(
        n_rows, n_cols,
        MARGIN_LEFT, table_top, CONTENT_WIDTH, table_height,
    )
    table = tbl_shape.table

    for col_idx, header_text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = header_text
        _set_font(run, size=TABLE_SIZE, bold=True, color=WHITE)
        _set_cell_fill(cell, ACCENT)
        _set_cell_margins(cell)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    for row_idx, row_data in enumerate(rows):
        bg_color = LIGHT_BG if row_idx % 2 == 0 else GRAY_BG
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = cell_text
            _set_font(run, size=TABLE_SIZE, bold=False, color=DARK_TEXT)
            _set_cell_fill(cell, bg_color)
            _set_cell_margins(cell)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    _set_table_borders(table, n_rows, n_cols)
    _add_footer(slide)


def add_section_slide(prs, title: str, subtitle: str = "") -> None:
    """챕터 전환 슬라이드. 오렌지 배경 + 큰 제목."""
    slide = _new_slide(prs)
    _add_rect(slide, Cm(0), Cm(0), SLIDE_WIDTH, SLIDE_HEIGHT, PRIMARY)

    _add_textbox(
        slide, Cm(3), Cm(6), SLIDE_WIDTH - Cm(6), Cm(4),
        text=title, font_size=Pt(36), bold=True, color=WHITE,
        alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )

    if subtitle:
        _add_textbox(
            slide, Cm(3), Cm(10.5), SLIDE_WIDTH - Cm(6), Cm(2),
            text=subtitle, font_size=SUBTITLE_SIZE, bold=False, color=WHITE,
            alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP,
        )

    _add_footer(slide)


def add_closing_slide(prs, title: str = "Thank You",
                      subtitle: str = "",
                      contact_lines: list[str] | None = None) -> None:
    """마무리 슬라이드. DARK_BG 배경 + 감사/연락처."""
    slide = _new_slide(prs)
    _add_rect(slide, Cm(0), Cm(0), SLIDE_WIDTH, SLIDE_HEIGHT, DARK_BG)

    # 오렌지 액센트 라인
    _add_rect(slide, Cm(12), Cm(8.5), Cm(10), Cm(0.08), PRIMARY)

    _add_textbox(
        slide, Cm(3), Cm(5), SLIDE_WIDTH - Cm(6), Cm(4),
        text=title, font_size=Pt(40), bold=True, color=WHITE,
        alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM,
    )

    if subtitle:
        _add_textbox(
            slide, Cm(3), Cm(9.5), SLIDE_WIDTH - Cm(6), Cm(2),
            text=subtitle, font_size=SUBTITLE_SIZE, bold=False, color=LIGHT_BG,
            alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP,
        )

    if contact_lines:
        y = Cm(12.5)
        for line in contact_lines:
            _add_textbox(
                slide, Cm(3), y, SLIDE_WIDTH - Cm(6), Cm(1),
                text=line, font_size=SMALL_SIZE, bold=False, color=LIGHT_BG,
                alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP,
            )
            y += Cm(0.8)

    _add_footer(slide)


# ---------------------------------------------------------------------------
# Additional slide helpers used by build_presentation.py
# ---------------------------------------------------------------------------

def add_screenshot_slide(prs, title: str, bullets: list[str],
                         screenshot_path: str,
                         screenshots_dir: Path | None = None) -> None:
    """스크린샷 슬라이드 (add_content_slide 래퍼)."""
    add_content_slide(prs, title, bullets, screenshot_path, screenshots_dir)


def add_diagram_slide(prs, title: str, text_content: str,
                      notes: list[str] | None = None) -> None:
    """텍스트 기반 다이어그램/파이프라인 슬라이드."""
    slide = _new_slide(prs)
    _add_header_bar(slide, title)

    content_height = SLIDE_HEIGHT - CONTENT_TOP - FOOTER_HEIGHT - Cm(0.5)

    if notes:
        diagram_height = content_height * 0.55
    else:
        diagram_height = content_height

    _add_rect(
        slide, MARGIN_LEFT + Cm(0.5), CONTENT_TOP + Cm(0.2),
        CONTENT_WIDTH - Cm(1), diagram_height, GRAY_BG,
    )

    _add_textbox(
        slide, MARGIN_LEFT + Cm(1), CONTENT_TOP + Cm(0.5),
        CONTENT_WIDTH - Cm(2), diagram_height - Cm(0.6),
        text=text_content, font_size=Pt(13), bold=False, color=DARK_TEXT,
        alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
    )

    if notes:
        notes_top = CONTENT_TOP + diagram_height + Cm(0.5)
        notes_height = content_height * 0.40
        txbox_n = _add_textbox(
            slide, MARGIN_LEFT, notes_top, CONTENT_WIDTH, notes_height,
        )
        _add_bullets_to_frame(
            txbox_n.text_frame, notes, font_size=Pt(13), color=DARK_TEXT,
        )

    _add_footer(slide)


def add_diagram_image_slide(prs, title: str, image_path: str,
                            notes: list[str] | None = None) -> None:
    """이미지 기반 다이어그램 슬라이드. HTML에서 캡처한 PNG를 삽입한다."""
    slide = _new_slide(prs)
    _add_header_bar(slide, title)

    content_height = SLIDE_HEIGHT - CONTENT_TOP - FOOTER_HEIGHT - Cm(0.5)

    if notes:
        diagram_height = content_height * 0.65
        notes_height = content_height * 0.30
    else:
        diagram_height = content_height

    try:
        slide.shapes.add_picture(
            image_path,
            MARGIN_LEFT + Cm(0.5), CONTENT_TOP + Cm(0.2),
            CONTENT_WIDTH - Cm(1), diagram_height,
        )
    except Exception as exc:
        print(f"    [WARN] 다이어그램 이미지 삽입 실패 ({image_path}): {exc}")
        _add_rect(
            slide, MARGIN_LEFT + Cm(0.5), CONTENT_TOP + Cm(0.2),
            CONTENT_WIDTH - Cm(1), diagram_height, GRAY_BG,
        )
        _add_textbox(
            slide, MARGIN_LEFT + Cm(1), CONTENT_TOP + Cm(3),
            CONTENT_WIDTH - Cm(2), Cm(2),
            text=f"[Diagram: {image_path}]",
            font_size=SMALL_SIZE, color=SECONDARY,
            alignment=PP_ALIGN.CENTER,
        )

    if notes:
        notes_top = CONTENT_TOP + diagram_height + Cm(0.5)
        txbox_n = _add_textbox(
            slide, MARGIN_LEFT, notes_top, CONTENT_WIDTH, notes_height,
        )
        _add_bullets_to_frame(
            txbox_n.text_frame, notes, font_size=Pt(13), color=DARK_TEXT,
        )

    _add_footer(slide)


def add_cards_slide(prs, title: str, cards: list[str], cols: int = 3) -> None:
    """카드 그리드 슬라이드."""
    slide = _new_slide(prs)
    _add_header_bar(slide, title)

    n_cards = len(cards)
    rows_count = (n_cards + cols - 1) // cols

    card_margin = Cm(0.4)
    total_avail_w = CONTENT_WIDTH - card_margin * (cols - 1)
    card_w = total_avail_w / cols

    content_avail_h = SLIDE_HEIGHT - CONTENT_TOP - FOOTER_HEIGHT - Cm(0.5)
    total_avail_h = content_avail_h - card_margin * (rows_count - 1)
    card_h = total_avail_h / rows_count

    for idx, card_text in enumerate(cards):
        r = idx // cols
        c = idx % cols
        left = MARGIN_LEFT + c * (card_w + card_margin)
        top = CONTENT_TOP + Cm(0.2) + r * (card_h + card_margin)

        bg_color = LIGHT_BG if idx % 2 == 0 else GRAY_BG
        _add_rect(slide, left, top, card_w, card_h, bg_color)

        if " \u2014 " in card_text:
            parts = card_text.split(" \u2014 ", 1)
            card_title = parts[0]
            card_desc = parts[1] if len(parts) > 1 else ""
        else:
            card_title = card_text
            card_desc = ""

        _add_textbox(
            slide, left + Cm(0.3), top + Cm(0.3),
            card_w - Cm(0.6), Cm(1.2),
            text=card_title, font_size=Pt(13), bold=True, color=PRIMARY,
            alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        )

        if card_desc:
            _add_textbox(
                slide, left + Cm(0.3), top + Cm(1.6),
                card_w - Cm(0.6), card_h - Cm(2.0),
                text=card_desc, font_size=Pt(11), bold=False, color=DARK_TEXT,
                alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            )

    _add_footer(slide)


def add_key_message_slide(prs, title: str, key_message: str,
                          sub_bullets: list[str]) -> None:
    """중앙 키 메시지 + 하위 불릿."""
    slide = _new_slide(prs)
    _add_header_bar(slide, title)

    _add_textbox(
        slide, Cm(2), CONTENT_TOP + Cm(1), SLIDE_WIDTH - Cm(4), Cm(3),
        text=key_message, font_size=Pt(22), bold=True, color=PRIMARY,
        alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )

    _add_rect(
        slide, Cm(8), CONTENT_TOP + Cm(4.5), SLIDE_WIDTH - Cm(16), Cm(0.05),
        ACCENT,
    )

    txbox = _add_textbox(
        slide, Cm(3), CONTENT_TOP + Cm(5.2), SLIDE_WIDTH - Cm(6),
        SLIDE_HEIGHT - CONTENT_TOP - Cm(5.2) - FOOTER_HEIGHT,
    )
    _add_bullets_to_frame(txbox.text_frame, sub_bullets, font_size=BODY_SIZE)

    _add_footer(slide)


def add_quote_slide(prs, title: str, quote: str, attribution: str,
                    feedback_points: list[str]) -> None:
    """인용문 스타일 슬라이드."""
    slide = _new_slide(prs)
    _add_header_bar(slide, title)

    quote_top = CONTENT_TOP + Cm(0.8)
    _add_rect(
        slide, Cm(2), quote_top, SLIDE_WIDTH - Cm(4), Cm(4.5), LIGHT_BG,
    )

    _add_textbox(
        slide, Cm(3), quote_top + Cm(0.5), SLIDE_WIDTH - Cm(6), Cm(2.5),
        text=quote, font_size=Pt(16), bold=False, color=DARK_TEXT,
        alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )

    _add_textbox(
        slide, Cm(3), quote_top + Cm(3), SLIDE_WIDTH - Cm(6), Cm(1),
        text=attribution, font_size=Pt(12), bold=False, color=SECONDARY,
        alignment=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.TOP,
    )

    feedback_top = quote_top + Cm(5.5)
    txbox = _add_textbox(
        slide, MARGIN_LEFT, feedback_top, CONTENT_WIDTH,
        SLIDE_HEIGHT - feedback_top - FOOTER_HEIGHT,
    )
    _add_bullets_to_frame(txbox.text_frame, feedback_points, font_size=BODY_SIZE)

    _add_footer(slide)


def add_three_column_slide(prs, title: str,
                           col_data: list[tuple[str, list[str]]]) -> None:
    """3단 컬럼 슬라이드."""
    slide = _new_slide(prs)
    _add_header_bar(slide, title)

    n_cols = len(col_data)
    col_gap = Cm(0.6)
    total_gap = col_gap * (n_cols - 1)
    col_width = (CONTENT_WIDTH - total_gap) / n_cols
    content_height = SLIDE_HEIGHT - CONTENT_TOP - FOOTER_HEIGHT - Cm(0.5)

    for i, (col_title, col_bullets) in enumerate(col_data):
        left = MARGIN_LEFT + i * (col_width + col_gap)

        bg_color = LIGHT_BG if i % 2 == 0 else GRAY_BG
        _add_rect(slide, left, CONTENT_TOP + Cm(0.2), col_width, content_height, bg_color)

        _add_textbox(
            slide, left + Cm(0.3), CONTENT_TOP + Cm(0.5),
            col_width - Cm(0.6), Cm(1.2),
            text=col_title, font_size=Pt(15), bold=True, color=PRIMARY,
            alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP,
        )

        txbox = _add_textbox(
            slide, left + Cm(0.3), CONTENT_TOP + Cm(2),
            col_width - Cm(0.6), content_height - Cm(2.5),
        )
        _add_bullets_to_frame(
            txbox.text_frame, col_bullets, font_size=Pt(12), color=DARK_TEXT,
        )

    _add_footer(slide)


def add_code_slide(prs, title: str, bullets: list[str],
                   code_text: str) -> None:
    """코드 블록 포함 슬라이드."""
    slide = _new_slide(prs)
    _add_header_bar(slide, title)

    content_height = SLIDE_HEIGHT - CONTENT_TOP - FOOTER_HEIGHT - Cm(0.5)

    txbox = _add_textbox(
        slide, MARGIN_LEFT, CONTENT_TOP, CONTENT_WIDTH, Cm(3),
    )
    _add_bullets_to_frame(txbox.text_frame, bullets, font_size=BODY_SIZE)

    code_top = CONTENT_TOP + Cm(3.5)
    code_height = content_height - Cm(3.8)
    _add_rect(
        slide, MARGIN_LEFT + Cm(0.3), code_top,
        CONTENT_WIDTH - Cm(0.6), code_height,
        CODE_BG,
    )

    _add_textbox(
        slide, MARGIN_LEFT + Cm(0.8), code_top + Cm(0.3),
        CONTENT_WIDTH - Cm(1.6), code_height - Cm(0.6),
        text=code_text, font_size=Pt(12), bold=False,
        color=CODE_TEXT,
        alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
    )

    _add_footer(slide)


def add_timeline_slide(prs, title: str, phases: list[tuple[str, str]]) -> None:
    """타임라인 슬라이드."""
    slide = _new_slide(prs)
    _add_header_bar(slide, title)

    n_phases = len(phases)
    content_height = SLIDE_HEIGHT - CONTENT_TOP - FOOTER_HEIGHT - Cm(1)
    phase_height = content_height / n_phases
    phase_gap = Cm(0.4)

    # 오렌지 그라데이션 순서 색상
    phase_colors = [PRIMARY, SECONDARY, ACCENT]

    for i, (phase_label, phase_desc) in enumerate(phases):
        top = CONTENT_TOP + Cm(0.5) + i * phase_height

        label_width = Cm(8)
        bg_color = phase_colors[i % len(phase_colors)]
        _add_rect(slide, MARGIN_LEFT, top, label_width, phase_height - phase_gap, bg_color)

        _add_textbox(
            slide, MARGIN_LEFT + Cm(0.3), top,
            label_width - Cm(0.6), phase_height - phase_gap,
            text=phase_label, font_size=Pt(14), bold=True, color=WHITE,
            alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
        )

        _add_textbox(
            slide, MARGIN_LEFT + label_width, top,
            Cm(1.5), phase_height - phase_gap,
            text="\u25B6", font_size=Pt(20), bold=False, color=bg_color,
            alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )

        desc_left = MARGIN_LEFT + label_width + Cm(1.5)
        desc_width = CONTENT_WIDTH - label_width - Cm(1.5)
        _add_rect(slide, desc_left, top, desc_width, phase_height - phase_gap, LIGHT_BG)
        _add_textbox(
            slide, desc_left + Cm(0.3), top,
            desc_width - Cm(0.6), phase_height - phase_gap,
            text=phase_desc, font_size=Pt(13), bold=False, color=DARK_TEXT,
            alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
        )

    _add_footer(slide)


# ---------------------------------------------------------------------------
# Sample PPT generation (직접 실행 시)
# ---------------------------------------------------------------------------

def build_sample_ppt():
    """7장짜리 샘플 PPT를 생성한다."""
    reset_slide_counter()
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # 1. Cover
    add_cover_slide(prs, "Secern AI", "프레젠테이션 템플릿",
                    extra_lines=["2026", "secern.ai"])

    # 2. Content
    add_content_slide(prs, "일반 콘텐츠 슬라이드", [
        "첫 번째 항목 설명",
        "두 번째 항목 설명",
        "세 번째 항목 설명",
    ])

    # 3. Two Column
    add_two_column_slide(prs, "비교 슬라이드",
        left_bullets=["왼쪽 항목 1", "왼쪽 항목 2", "왼쪽 항목 3"],
        right_bullets=["오른쪽 항목 1", "오른쪽 항목 2", "오른쪽 항목 3"],
        left_title="옵션 A", right_title="옵션 B",
    )

    # 4. Table
    add_table_slide(prs, "데이터 테이블 슬라이드",
        headers=["항목", "값", "상태"],
        rows=[
            ["항목 1", "100", "완료"],
            ["항목 2", "200", "진행 중"],
            ["항목 3", "300", "예정"],
        ],
    )

    # 5. Section
    add_section_slide(prs, "다음 섹션", "챕터 전환 슬라이드 예시")

    # 6. Code
    add_code_slide(prs, "코드 블록 슬라이드", [
        "API 호출 예시",
    ], 'curl -X POST https://api.secern.ai/generate \\\n  -H "Authorization: Bearer $TOKEN" \\\n  -d \'{"prompt": "hello"}\'')

    # 7. Closing
    add_closing_slide(prs, "Thank You", "Secern AI",
                      contact_lines=["secern.ai", "info@secern.ai"])

    output = Path(__file__).resolve().parent / "Secern_AI_Template_Sample.pptx"
    prs.save(str(output))
    print(f"샘플 PPT 생성 완료: {output}")
    print(f"  슬라이드 수: {_slide_number}")
    print(f"  파일 크기: {output.stat().st_size / 1024:.1f} KB")
    return output


if __name__ == "__main__":
    build_sample_ppt()
